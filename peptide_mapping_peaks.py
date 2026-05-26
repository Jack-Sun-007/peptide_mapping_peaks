import re
import os
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
import PySimpleGUI as sg


# 保留时间按比例计算到ppt尺寸
def get_RT(a, start, end):
    b = ((float(a)-start-3.1)/end*36.2)
    return Cm(b)


# 使用正则表达式去除氨基酸的大写字母
def remove_uppercase(s):
    a = re.sub(r'([A-Za-z0-9]+):[A-Za-z]+(\d+)-[A-Za-z]+(\d+)', r'\1:\2-\3', s)
    return a


def main(data, start_time, end_time, text_size, show_time, remove_upper):
    # 创建新的 PowerPoint 演示文稿
    ppt = Presentation()
    # 设置幻灯片大小为16:9
    ppt.slide_width = Cm(33.867)
    ppt.slide_height = Cm(19.05)
    # 添加一个新幻灯片
    slide_layout = ppt.slide_layouts[6]  # 空白幻灯片布局
    slide = ppt.slides.add_slide(slide_layout)
    # 定义初始位置
    left = Cm(0)
    top = Cm(4)
    width = Cm(2.8)
    height = Cm(0.77)
    # 处理文本数据，按行分割并将每个条目各自作为一个文本框
    raw = data.split('\n')
    for line in raw:
        parts = line.strip().split('\t')  # 以制表符为分隔符，获取每一部分
        part1 = parts[0]
        part2 = parts[1]
        RT = get_RT(part2, int(start_time), int(end_time))
        # 添加第一个文本框，肽段名称
        textbox1 = slide.shapes.add_textbox(RT, top, width, height)
        text_frame1 = textbox1.text_frame
        if remove_upper == 1:
            text_frame1.text = remove_uppercase(part1)  # 对肽段名称去掉字母
        else:
            text_frame1.text = part1
        p1 = text_frame1.paragraphs[0]
        p1.font.name = 'Times New Roman'
        p1.font.size = Pt(text_size)
        p1.font.color.rgb = RGBColor(68, 114, 196)  # 深蓝
        # 旋转文本框
        textbox1.rotation = -90
        # 添加第二个文本框，保留时间
        if show_time is True:
            textbox2 = slide.shapes.add_textbox(RT, top - Cm(4), width, height)
            text_frame2 = textbox2.text_frame
            text_frame2.text = part2  # 直接设置内容
            p2 = text_frame2.paragraphs[0]
            p2.font.name = 'Times New Roman'
            p2.font.size = Pt(text_size)
            p2.font.color.rgb = RGBColor(192, 0, 0)  # 深红
            # 旋转文本框
            textbox2.rotation = -90
        else:
            pass
    # 保存 PowerPoint 文件
    ppt.save('PM_peak.pptx')
    os.startfile('PM_peak.pptx')


def GUI():
    sg.theme('GrayGrayGray')
    # 界面布局，将会按照列表顺序从上往下依次排列，二级列表中，从左往右依此排列
    layout = [[sg.Text('请输入标峰字号大小'), sg.InputText('11', size=5)],
              [sg.Checkbox('显示保留时间', default=True), sg.Checkbox('去除氨基酸字母', default=True)],
              [sg.Text('请输入标峰开始的保留时间(min)'), sg.InputText('5', size=5)],
              [sg.Text('请输入标峰结束的保留时间(min)'), sg.InputText('80', size=5)],
              [sg.Text('请复制标峰肽段和保留时间：')],
              [sg.Multiline('LC:S208-R211\t12.30\nLC:V104-R108\t16.63\nLC:A25-K42\t37.11', size=(40, 40)), sg.Button('运行')]
              ]

    # 创造窗口
    window = sg.Window('肽图&二硫键标峰辅助', layout, size=(400, 300))
    # 事件循环并获取输入值
    while True:
        event, values = window.read()
        if event in (None, 'Cancel'):  # 如果用户关闭窗口或点击`Cancel`
            break
        #print(values)
        text_size = int(values[0])
        show_time = values[1]
        remove_upper = values[2]
        start_time = int(values[3])
        end_time = int(values[4])
        data = values[5]
        main(data, start_time, end_time, text_size, show_time, remove_upper)
        break
    window.close()


if __name__ == '__main__':
    try:
        GUI()
    except Exception as e:
        sg.popup_error(f'运行出错!\n报错信息：', e)
